"""
make_presentation.py
====================

Build a 6-slide PowerPoint deck summarising the XGBoost ignition model
trained on ``database_xgb.csv`` (see ``xgb_ignition_model_2.py``).

Slide 1  - Architecture and design justifications.
Slide 2  - Headline performance (hold-out + CV).
Slide 3  - ROC and Precision-Recall curves.
Slide 4  - Confusion matrices and threshold tuning.
Slide 5  - Feature importance: physics vs memorisation.
Slide 6  - The stratified-vs-group-CV gap (true generalisation).

Run::

    python make_presentation.py
        --> writes presentation/xgb_ignition_v2_deck.pptx
        --> writes presentation/figures/*.png
"""

from __future__ import annotations

import json
from pathlib import Path

import joblib
import matplotlib

matplotlib.use("Agg")
import matplotlib.patches as mpatches
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt
from sklearn.metrics import (
    average_precision_score,
    confusion_matrix,
    precision_recall_curve,
    roc_auc_score,
    roc_curve,
)
from sklearn.model_selection import (
    GroupKFold,
    StratifiedKFold,
    cross_val_score,
    train_test_split,
)

from xgb_ignition_model_2 import (
    CATEGORICAL_FEATURES,
    NUMERIC_FEATURES,
    _onehot_feature_names,
    load_clean,
)

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

DATA_PATH = Path("database_xgb.csv")
ARTIFACTS = Path("artifacts_v2")
OUT_DIR = Path("presentation")
FIG_DIR = OUT_DIR / "figures"
OUT_DIR.mkdir(parents=True, exist_ok=True)
FIG_DIR.mkdir(parents=True, exist_ok=True)

RANDOM_STATE = 42

# Colour palette
COL_PRIMARY = "#1f3a5f"
COL_SECONDARY = "#c25450"
COL_ACCENT = "#3d8a8a"
COL_LIGHT = "#eef2f6"
COL_TEXT = "#1f1f1f"
COL_GREY = "#7a7a7a"

plt.rcParams.update(
    {
        "font.family": "DejaVu Sans",
        "font.size": 12,
        "axes.titlesize": 14,
        "axes.labelsize": 12,
        "axes.edgecolor": "#333333",
        "axes.linewidth": 1.0,
        "axes.spines.top": False,
        "axes.spines.right": False,
        "xtick.color": "#333333",
        "ytick.color": "#333333",
        "grid.color": "#dddddd",
        "grid.linestyle": "--",
        "grid.linewidth": 0.6,
        "figure.dpi": 150,
    }
)

# ---------------------------------------------------------------------------
# 1. Re-run minimal evaluation to produce fresh figures
# ---------------------------------------------------------------------------

print("[1/4] Loading data + model ...")
df = load_clean(DATA_PATH)
X = df[NUMERIC_FEATURES + CATEGORICAL_FEATURES]
y = df["ignition_binary"].astype(int)
groups = df["source_group"]

pipe = joblib.load(ARTIFACTS / "xgb_ignition_model_v2.joblib")

# Rebuild the same 80/20 stratified split used during training so that
# hold-out figures are identical to the metrics in ``metrics.json``.
X_train, X_test, y_train, y_test = train_test_split(
    X, y, test_size=0.20, stratify=y, random_state=RANDOM_STATE
)

proba_test = pipe.predict_proba(X_test)[:, 1]
roc_auc = roc_auc_score(y_test, proba_test)
pr_auc = average_precision_score(y_test, proba_test)

with open(ARTIFACTS / "metrics.json") as f:
    metrics = json.load(f)

print(f"  -> rebuilt hold-out ROC-AUC: {roc_auc:.4f}, PR-AUC: {pr_auc:.4f}")

# ---------------------------------------------------------------------------
# 2. Figure helpers
# ---------------------------------------------------------------------------

def _save_fig(fig, name: str) -> Path:
    p = FIG_DIR / name
    fig.tight_layout()
    fig.savefig(p, dpi=180, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return p


# --- Architecture diagram (Slide 1) ---------------------------------------

def fig_architecture() -> Path:
    fig, ax = plt.subplots(figsize=(11, 5.5))
    ax.axis("off")
    ax.set_xlim(0, 11)
    ax.set_ylim(0, 6)

    def box(x, y, w, h, text, fill=COL_LIGHT, edge=COL_PRIMARY, fc_text=COL_TEXT, fs=11, weight="bold"):
        rect = mpatches.FancyBboxPatch(
            (x, y), w, h,
            boxstyle="round,pad=0.05,rounding_size=0.12",
            linewidth=1.5, edgecolor=edge, facecolor=fill,
        )
        ax.add_patch(rect)
        ax.text(x + w / 2, y + h / 2, text, ha="center", va="center",
                fontsize=fs, color=fc_text, weight=weight, wrap=True)

    def arrow(x1, y1, x2, y2):
        ax.annotate("", xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle="->", lw=1.6, color=COL_PRIMARY))

    box(0.2, 4.6, 2.0, 0.9, "Raw CSV\ndatabase_xgb.csv\n5 208 rows", fill="#ffffff")

    box(0.2, 1.6, 2.0, 2.4,
        "Drop leakage\n* trailing 'Unnamed'\n* post-ignition cols\n  - Flame Length\n  - FSR / HRR\n  - Smoke / Aerosols",
        fill="#fff3f0", edge=COL_SECONDARY, fs=9, weight="normal")

    arrow(1.2, 4.6, 1.2, 4.05)
    arrow(2.2, 2.8, 2.85, 3.8)
    arrow(2.2, 5.05, 2.85, 4.8)

    box(2.9, 3.5, 2.7, 2.0,
        "Cleaning pipeline\n* SI units (kPa, mm/s, g)\n* Canonical categories\n  (lowercase merge, typos)\n* Engineer dims, energy",
        fill=COL_LIGHT, fs=10, weight="normal")

    arrow(5.6, 4.5, 6.25, 4.5)

    box(6.3, 3.5, 2.6, 2.0,
        "ColumnTransformer\n* Numeric: pass-through\n  (XGBoost NaN-native)\n* Categorical: one-hot\n  + 'Unknown' bin",
        fill=COL_LIGHT, fs=10, weight="normal")

    arrow(8.9, 4.5, 9.45, 4.5)

    box(9.5, 3.4, 1.4, 2.2,
        "XGBoost\nhist trees\nscale_pos_w\n= neg/pos\n= 0.347",
        fill=COL_ACCENT, edge=COL_PRIMARY, fc_text="white", fs=10)

    box(2.9, 0.2, 8.0, 1.2,
        "Validation\n* Stratified 80/20 hold-out   * Stratified 5-fold CV (in-distribution)\n* GroupKFold-by-DOI 5-fold CV (cross-source - honest)   * F1-optimal threshold from PR curve",
        fill="#f4f9f9", edge=COL_ACCENT, fs=10, weight="normal")

    arrow(10.2, 3.4, 8.8, 1.45)

    return _save_fig(fig, "architecture.png")


# --- Headline performance KPI panel (Slide 2) ------------------------------

def fig_kpi_panel() -> Path:
    fig, ax = plt.subplots(figsize=(11, 4.5))
    ax.axis("off")
    ax.set_xlim(0, 11)
    ax.set_ylim(0, 4)

    kpis = [
        ("Hold-out\nROC-AUC", f"{metrics['holdout_roc_auc']:.3f}", COL_PRIMARY),
        ("Hold-out\nPR-AUC", f"{metrics['holdout_pr_auc']:.3f}", COL_ACCENT),
        ("Best F1\n(tuned thr.)", f"{metrics['holdout_best_f1']:.3f}", COL_SECONDARY),
        ("Stratified\n5-fold ROC-AUC", f"{metrics['stratified_5fold_roc_auc_mean']:.3f}\n+/- {metrics['stratified_5fold_roc_auc_std']:.3f}", COL_PRIMARY),
        ("Group-by-DOI\n5-fold ROC-AUC", f"{metrics['groupkfold_5fold_roc_auc_mean']:.3f}\n+/- {metrics['groupkfold_5fold_roc_auc_std']:.2f}", COL_SECONDARY),
    ]
    w = 1.95
    gap = 0.15
    total = len(kpis) * w + (len(kpis) - 1) * gap
    x0 = (11 - total) / 2
    for i, (label, value, colour) in enumerate(kpis):
        x = x0 + i * (w + gap)
        rect = mpatches.FancyBboxPatch(
            (x, 0.5), w, 3.0,
            boxstyle="round,pad=0.04,rounding_size=0.18",
            linewidth=2.0, edgecolor=colour, facecolor="white",
        )
        ax.add_patch(rect)
        ax.text(x + w / 2, 2.7, value, ha="center", va="center",
                fontsize=20, color=colour, weight="bold")
        ax.text(x + w / 2, 1.1, label, ha="center", va="center",
                fontsize=11, color=COL_TEXT)
    return _save_fig(fig, "kpis.png")


# --- ROC + PR curves side by side (Slide 3) --------------------------------

def fig_roc_pr() -> Path:
    fig, axes = plt.subplots(1, 2, figsize=(11, 4.5))

    # ROC
    fpr, tpr, _ = roc_curve(y_test, proba_test)
    axes[0].plot(fpr, tpr, color=COL_PRIMARY, lw=2.4,
                 label=f"XGBoost (AUC = {roc_auc:.3f})")
    axes[0].plot([0, 1], [0, 1], "k--", alpha=0.4, label="Chance")
    axes[0].fill_between(fpr, tpr, alpha=0.08, color=COL_PRIMARY)
    axes[0].set_xlabel("False positive rate")
    axes[0].set_ylabel("True positive rate")
    axes[0].set_title("ROC curve (hold-out)")
    axes[0].set_xlim(0, 1)
    axes[0].set_ylim(0, 1.02)
    axes[0].grid(True, alpha=0.4)
    axes[0].legend(loc="lower right", frameon=False)

    # PR
    prec, rec, _ = precision_recall_curve(y_test, proba_test)
    baseline = float(y_test.mean())
    axes[1].plot(rec, prec, color=COL_ACCENT, lw=2.4,
                 label=f"XGBoost (AP = {pr_auc:.3f})")
    axes[1].axhline(baseline, ls="--", color="k", alpha=0.4,
                    label=f"Baseline = {baseline:.2f}")
    axes[1].fill_between(rec, prec, alpha=0.08, color=COL_ACCENT)
    axes[1].set_xlabel("Recall")
    axes[1].set_ylabel("Precision")
    axes[1].set_title("Precision-Recall curve (hold-out)")
    axes[1].set_xlim(0, 1)
    axes[1].set_ylim(0, 1.02)
    axes[1].grid(True, alpha=0.4)
    axes[1].legend(loc="lower left", frameon=False)

    return _save_fig(fig, "roc_pr.png")


# --- Confusion matrices at two thresholds (Slide 4) ------------------------

def _draw_cm(ax, cm, title):
    im = ax.imshow(cm, cmap="Blues", vmin=0, vmax=cm.max())
    for (i, j), v in np.ndenumerate(cm):
        ax.text(j, i, str(v), ha="center", va="center",
                color="white" if v > cm.max() / 2 else "black",
                fontsize=15, weight="bold")
    ax.set_xticks([0, 1], ["No ign.", "Ignition"])
    ax.set_yticks([0, 1], ["No ign.", "Ignition"])
    ax.set_xlabel("Predicted")
    ax.set_ylabel("True")
    ax.set_title(title)
    return im


def fig_confusion() -> Path:
    fig, axes = plt.subplots(1, 2, figsize=(10, 4.6))
    thr_default = 0.5
    thr_best = metrics["holdout_best_threshold"]

    pred_def = (proba_test >= thr_default).astype(int)
    pred_best = (proba_test >= thr_best).astype(int)

    cm1 = confusion_matrix(y_test, pred_def)
    cm2 = confusion_matrix(y_test, pred_best)

    _draw_cm(axes[0], cm1, f"Threshold = 0.50 (default)\nF1 = {metrics['holdout_f1_default_thr']:.3f}")
    _draw_cm(axes[1], cm2, f"Threshold = {thr_best:.2f} (F1-optimal)\nF1 = {metrics['holdout_best_f1']:.3f}")

    return _save_fig(fig, "confusion.png")


# --- Top-15 feature importance bar chart (Slide 5) -------------------------

def fig_feature_importance() -> Path:
    imp_df = pd.read_csv(ARTIFACTS / "feature_importance.csv")
    top = imp_df.head(15).iloc[::-1]  # reverse so largest is at top

    # Colour rig-related bars distinctly to highlight the memorisation story.
    colours = [
        COL_SECONDARY if f.startswith("rig_grouped=") else COL_PRIMARY
        for f in top["feature"]
    ]

    fig, ax = plt.subplots(figsize=(10, 5.5))
    ax.barh(top["feature"], top["importance"], color=colours, edgecolor="white")
    ax.set_xlabel("Gain-based importance")
    ax.set_title("Top 15 features by XGBoost gain")
    ax.grid(True, axis="x", alpha=0.4)
    ax.spines["left"].set_color("#cccccc")

    rig_patch = mpatches.Patch(color=COL_SECONDARY, label="Rig identity (one-hot)")
    other_patch = mpatches.Patch(color=COL_PRIMARY, label="Physical / chemistry features")
    ax.legend(handles=[rig_patch, other_patch], loc="lower right", frameon=False)

    return _save_fig(fig, "importance.png")


# --- Stratified vs Group CV comparison (Slide 6) ---------------------------

def fig_cv_gap() -> Path:
    # Compute per-fold scores so we can show distributions, not just means.
    print("  computing per-fold CV scores for slide 6 ...")
    skf = StratifiedKFold(n_splits=5, shuffle=True, random_state=RANDOM_STATE)
    skf_scores = cross_val_score(pipe, X, y, cv=skf, scoring="roc_auc", n_jobs=-1)
    gkf = GroupKFold(n_splits=5)
    gkf_scores = cross_val_score(
        pipe, X, y, cv=gkf, groups=groups, scoring="roc_auc", n_jobs=-1
    )

    fig, ax = plt.subplots(figsize=(10, 5))
    means = [skf_scores.mean(), gkf_scores.mean()]
    stds = [skf_scores.std(), gkf_scores.std()]
    xs = [0.4, 1.4]
    ax.bar(
        xs, means, width=0.55,
        yerr=stds, capsize=8,
        color=[COL_PRIMARY, COL_SECONDARY],
        edgecolor="white", linewidth=1.5,
    )
    # Scatter per-fold points
    for x, scores, c in [
        (0.4, skf_scores, COL_PRIMARY),
        (1.4, gkf_scores, COL_SECONDARY),
    ]:
        ax.scatter(
            np.full_like(scores, x) + np.random.RandomState(0).uniform(-0.06, 0.06, len(scores)),
            scores, s=70, color="white", edgecolor=c, linewidth=2, zorder=5,
        )

    for x, m, s in zip(xs, means, stds):
        ax.text(x, m + s + 0.025, f"{m:.3f}\n+/- {s:.3f}",
                ha="center", va="bottom", weight="bold", fontsize=11)

    ax.axhline(0.5, color="k", lw=0.8, ls="--", alpha=0.5)
    ax.text(1.85, 0.51, "Chance", color="k", alpha=0.6, fontsize=9)

    ax.set_xticks(xs, [
        "Stratified 5-fold CV\n(rows shuffled)",
        "GroupKFold 5-fold CV\n(by source DOI)",
    ])
    ax.set_ylabel("ROC-AUC (5-fold CV)")
    ax.set_title("In-distribution vs cross-source generalisation")
    ax.set_ylim(0.4, 1.0)
    ax.grid(True, axis="y", alpha=0.4)

    gap = skf_scores.mean() - gkf_scores.mean()
    ax.annotate(
        f"gap = {gap:.3f}\n(source leakage)",
        xy=(1.0, (skf_scores.mean() + gkf_scores.mean()) / 2),
        ha="center", va="center", fontsize=11, color=COL_TEXT,
        bbox=dict(boxstyle="round,pad=0.4", fc="#fff5e6", ec="#e6a23c"),
    )

    return _save_fig(fig, "cv_gap.png"), skf_scores, gkf_scores


# ---------------------------------------------------------------------------
# 3. Generate all figures
# ---------------------------------------------------------------------------

print("[2/4] Generating figures ...")
fig_arch_path = fig_architecture()
fig_kpi_path = fig_kpi_panel()
fig_rocpr_path = fig_roc_pr()
fig_cm_path = fig_confusion()
fig_imp_path = fig_feature_importance()
fig_cvgap_path, skf_scores, gkf_scores = fig_cv_gap()

# ---------------------------------------------------------------------------
# 4. Build the PowerPoint
# ---------------------------------------------------------------------------

print("[3/4] Building PowerPoint ...")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

blank_layout = prs.slide_layouts[6]


def _add_rect(slide, x, y, w, h, fill_rgb, line_rgb=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor.from_string(fill_rgb.replace("#", ""))
    shp.line.fill.background() if line_rgb is None else _set_line(shp, line_rgb)
    shp.shadow.inherit = False
    return shp


def _set_line(shp, hexcolor):
    shp.line.color.rgb = RGBColor.from_string(hexcolor.replace("#", ""))


def _add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=COL_TEXT, align="left"):
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = "Calibri"
        run.font.color.rgb = RGBColor.from_string(color.replace("#", ""))
    return tb


def _add_bullets(slide, x, y, w, h, items, *, size=14, color=COL_TEXT, bold_first=False):
    """Items: list of strings, optional leading '* ' is rendered as a bullet."""
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.name = "Calibri"
        run.font.color.rgb = RGBColor.from_string(color.replace("#", ""))
        if bold_first and i == 0:
            run.font.bold = True
    return tb


def _add_title_bar(slide, title, subtitle=None):
    _add_rect(slide, Inches(0), Inches(0), SW, Inches(0.85), COL_PRIMARY)
    _add_rect(slide, Inches(0), Inches(0.85), SW, Inches(0.06), COL_SECONDARY)
    _add_text(
        slide, Inches(0.4), Inches(0.12), SW - Inches(0.8), Inches(0.6),
        title, size=24, bold=True, color="#ffffff",
    )
    if subtitle:
        _add_text(
            slide, Inches(0.4), Inches(0.55), SW - Inches(0.8), Inches(0.35),
            subtitle, size=12, color="#dde6f0",
        )


def _add_footer(slide, page, total=6):
    _add_text(
        slide, Inches(0.4), SH - Inches(0.35), Inches(8), Inches(0.3),
        "XGBoost ignition classifier - database_xgb.csv (4 449 clean rows)",
        size=9, color=COL_GREY,
    )
    _add_text(
        slide, SW - Inches(1.3), SH - Inches(0.35), Inches(1.0), Inches(0.3),
        f"{page} / {total}", size=9, color=COL_GREY, align="right",
    )


def _add_image(slide, path, x, y, w=None, h=None):
    if w is not None and h is not None:
        slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    elif w is not None:
        slide.shapes.add_picture(str(path), x, y, width=w)
    elif h is not None:
        slide.shapes.add_picture(str(path), x, y, height=h)
    else:
        slide.shapes.add_picture(str(path), x, y)


# -------- Slide 1: Architecture --------------------------------------------

s1 = prs.slides.add_slide(blank_layout)
_add_title_bar(
    s1,
    "Model architecture and design choices",
    "Why XGBoost on a literature-aggregated combustion database",
)
_add_image(s1, fig_arch_path, Inches(0.35), Inches(1.1), w=Inches(7.6))

_add_text(s1, Inches(8.2), Inches(1.1), Inches(4.9), Inches(0.4),
          "Justifications", size=15, bold=True, color=COL_PRIMARY)
_add_bullets(
    s1, Inches(8.2), Inches(1.55), Inches(4.95), Inches(5.3),
    [
        "* Tabular + non-linear interactions (O2, pressure,",
        "  gravity, flow, material) -> gradient-boosted trees",
        "  dominate this regime.",
        "* Native NaN handling lets predictors that are",
        "  physically inapplicable (e.g. core_diameter on a flat",
        "  sample) be 'missing' = information, not noise.",
        "* Tree splits are scale-invariant -> we only need",
        "  consistent units across rows, no log / standardise.",
        "* L1/L2 + row/col subsampling + scale_pos_weight",
        "  = neg/pos = 0.347 controls overfitting on ~4.5k rows",
        "  while focusing gradient on the minority no-ignition",
        "  class (24 % of data).",
        "* Post-ignition columns (Flame Length, FSR, HRR,",
        "  Smoke) dropped -> no target leakage.",
        "* DOI kept as a grouping key so we can validate on",
        "  unseen experimental campaigns, not just unseen rows.",
        "* Cleaning collapses unit / case / typo variants:",
        "  '21%' = '0.21', 'Flat' = 'flat',",
        "  'Open Flame' = 'open flame'.",
    ],
    size=12,
)
_add_footer(s1, 1)

# -------- Slide 2: Headline performance ------------------------------------

s2 = prs.slides.add_slide(blank_layout)
_add_title_bar(
    s2,
    "Headline performance",
    "80/20 stratified hold-out + 5-fold cross-validation",
)
_add_image(s2, fig_kpi_path, Inches(0.7), Inches(1.1), w=Inches(11.9))

_add_text(s2, Inches(0.6), Inches(4.5), Inches(12), Inches(0.4),
          "Reading the numbers", size=15, bold=True, color=COL_PRIMARY)
_add_bullets(
    s2, Inches(0.6), Inches(4.95), Inches(12.2), Inches(2.3),
    [
        "* ROC-AUC = 0.90 on a stratified hold-out -> for ~9 out of 10 randomly chosen (ignition, no-ignition) pairs the model",
        "  scores the igniting case higher. Stratified 5-fold CV confirms it (0.903 +/- 0.004), so the hold-out result is not lucky.",
        "* PR-AUC = 0.96 is the more honest figure under our 74/26 class imbalance: even when we slide along the precision/recall",
        "  curve, average precision stays near-perfect because true ignitions are densely concentrated at high model scores.",
        "* Default-threshold F1 = 0.87 jumps to 0.89 once we move the decision boundary from 0.5 to 0.25, which is exactly what",
        "  we expect after down-weighting the majority class (scale_pos_weight = 0.347 deflates positive scores).",
        "* Over v1 (database.csv) every metric improves and the stratified-CV std halves, consistent with ~60 % more clean rows.",
    ],
    size=12,
)
_add_footer(s2, 2)

# -------- Slide 3: ROC + PR -------------------------------------------------

s3 = prs.slides.add_slide(blank_layout)
_add_title_bar(
    s3,
    "Ranking quality: ROC and Precision-Recall curves",
    "How discriminative is the predicted probability?",
)
_add_image(s3, fig_rocpr_path, Inches(0.4), Inches(1.05), w=Inches(8.5))
_add_text(s3, Inches(9.1), Inches(1.05), Inches(4.0), Inches(0.4),
          "Interpretation", size=15, bold=True, color=COL_PRIMARY)
_add_bullets(
    s3, Inches(9.1), Inches(1.5), Inches(4.05), Inches(5.5),
    [
        "ROC curve",
        "* Hugs the top-left corner: at 80 % recall",
        "  on ignition the false-positive rate on",
        "  no-ignition stays under ~20 %.",
        "* The curve is symmetric: the model is",
        "  equally good at ranking ignitions up and",
        "  no-ignitions down - no class is being",
        "  ignored.",
        "",
        "PR curve",
        "* Average precision = 0.96, well above the",
        "  baseline 0.74 (positive prevalence).",
        "* Precision stays >= 0.90 up to recall ~0.90,",
        "  which is the operating regime we care",
        "  about: avoid false 'ignites' in safety-",
        "  critical screening.",
        "* The PR view is the right one here because",
        "  ROC-AUC over-rewards rejecting easy",
        "  negatives when the positive class is the",
        "  majority - PR is the realistic metric.",
    ],
    size=11,
)
_add_footer(s3, 3)

# -------- Slide 4: Confusion + threshold ------------------------------------

s4 = prs.slides.add_slide(blank_layout)
_add_title_bar(
    s4,
    "Operating point: confusion matrices and threshold tuning",
    "Default 0.5 vs F1-optimal 0.25 -- precision/recall trade-off",
)
_add_image(s4, fig_cm_path, Inches(0.4), Inches(1.05), w=Inches(8.4))

thr_best = metrics["holdout_best_threshold"]
_add_text(s4, Inches(9.0), Inches(1.05), Inches(4.1), Inches(0.4),
          "Why two thresholds?", size=15, bold=True, color=COL_PRIMARY)
_add_bullets(
    s4, Inches(9.0), Inches(1.5), Inches(4.1), Inches(5.5),
    [
        "Threshold 0.50 (default)",
        "* Recall on no-ignition = 0.83 (190/229),",
        "  recall on ignition = 0.81 (537/661).",
        "* Accuracy 82 %, F1 = 0.87.",
        "* Bias caused by scale_pos_weight: the",
        "  model deflates positive scores, so 0.5 is",
        "  too conservative for the majority class.",
        "",
        f"Threshold {thr_best:.2f} (F1-optimal)",
        "* Ignition recall jumps to 0.92 (605/661);",
        "  no-ignition recall drops to 0.60 (137/229).",
        "* Best F1 = 0.89, accuracy 83 %.",
        "",
        "Pick the threshold to your cost:",
        "* Safety screening: keep a low threshold ->",
        "  high ignition recall, accept some false",
        "  positives.",
        "* Material-down-selection: raise the",
        "  threshold to maximise precision on the",
        "  'no-ignition' (safe) class.",
    ],
    size=11,
)
_add_footer(s4, 4)

# -------- Slide 5: Feature importance ---------------------------------------

s5 = prs.slides.add_slide(blank_layout)
_add_title_bar(
    s5,
    "Feature importance: physics vs memorisation",
    "Gain-based importances on the one-hot expanded feature space",
)
_add_image(s5, fig_imp_path, Inches(0.4), Inches(1.05), w=Inches(8.4))

_add_text(s5, Inches(9.0), Inches(1.05), Inches(4.1), Inches(0.4),
          "What the model is using", size=15, bold=True, color=COL_PRIMARY)
_add_bullets(
    s5, Inches(9.0), Inches(1.5), Inches(4.1), Inches(5.7),
    [
        "Rig identity dominates (red bars)",
        "* FLARE, GIFFTS, HOI test apparatus, SJ-10",
        "  satellite, NASA Glenn ZGF together",
        "  account for ~25 % of total gain.",
        "* This is partly real (each rig has its own",
        "  flow / radiative environment) and partly",
        "  memorisation of campaign-specific",
        "  conditions -> see slide 6.",
        "",
        "Physical drivers (blue bars)",
        "* Material family: Kimwipes (cellulosic),",
        "  PMMA, PDMS, Kapton, NiCr/ETFE wire",
        "  composites - matches combustion intuition.",
        "* Geometry = Flat is a strong split (flat",
        "  samples behave differently than wires).",
        "* outer_diameter_mm, sample_dim_mean_mm,",
        "  ignition_time_s - all are first-order",
        "  ignition-energy proxies.",
        "* flow_direction = Quiescent captures the",
        "  no-flow micro-g regime where opposed-",
        "  flow ignition rules don't apply.",
    ],
    size=11,
)
_add_footer(s5, 5)

# -------- Slide 6: Stratified vs Group CV gap -------------------------------

s6 = prs.slides.add_slide(blank_layout)
_add_title_bar(
    s6,
    "The honest generalisation gap",
    "Random folds (0.90) vs Group-by-DOI folds (0.68)",
)
_add_image(s6, fig_cvgap_path, Inches(0.4), Inches(1.05), w=Inches(8.4))

gap = skf_scores.mean() - gkf_scores.mean()
_add_text(s6, Inches(9.0), Inches(1.05), Inches(4.1), Inches(0.4),
          "Why the gap matters", size=15, bold=True, color=COL_PRIMARY)
_add_bullets(
    s6, Inches(9.0), Inches(1.5), Inches(4.1), Inches(5.7),
    [
        f"Gap = {gap:.2f} ROC-AUC",
        "* Stratified random CV mixes rows from the",
        "  same paper across train and test folds.",
        "  XGBoost can latch onto rig-specific",
        "  fingerprints (rare flow speeds, fixed",
        "  ignition powers, characteristic material",
        "  cocktails) and 'memorise' them.",
        "* GroupKFold holds out entire papers/DOIs:",
        "  this rig-fingerprint shortcut disappears,",
        "  exposing how much true physics the model",
        "  learnt.",
        "",
        "What to quote",
        "* For a NEW row from a KNOWN rig: 0.90",
        "  ROC-AUC is realistic.",
        "* For a NEW rig / NEW campaign: 0.68 +/-",
        "  0.10 ROC-AUC is the honest expectation.",
        "",
        "Paths to close the gap",
        "* Collect more diverse sources (more rigs).",
        "* Per-rig calibration / mixed-effects layer.",
        "* Domain-adversarial training to make",
        "  representations rig-invariant.",
    ],
    size=11,
)
_add_footer(s6, 6)

out_path = OUT_DIR / "xgb_ignition_v2_deck.pptx"
prs.save(str(out_path))

print(f"[4/4] Saved -> {out_path.resolve()}")
